from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

# Create presentation
prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(7.5)

# Define color scheme
DARK_BLUE = RGBColor(25, 56, 109)
LIGHT_BLUE = RGBColor(65, 105, 225)
ACCENT_ORANGE = RGBColor(255, 140, 0)
WHITE = RGBColor(255, 255, 255)
DARK_GRAY = RGBColor(64, 64, 64)
LIGHT_GRAY = RGBColor(240, 240, 240)

def add_title_slide(prs, title, subtitle=""):
    """Add a title slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = DARK_BLUE
    
    # Add title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(1.5))
    title_frame = title_box.text_frame
    title_frame.word_wrap = True
    p = title_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(60)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
    
    # Add subtitle
    if subtitle:
        subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.2), Inches(9), Inches(1.5))
        subtitle_frame = subtitle_box.text_frame
        subtitle_frame.word_wrap = True
        p = subtitle_frame.paragraphs[0]
        p.text = subtitle
        p.font.size = Pt(28)
        p.font.color.rgb = ACCENT_ORANGE
        p.alignment = PP_ALIGN.CENTER
    
    return slide

def add_content_slide(prs, title, content_items):
    """Add a content slide with bullet points"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = WHITE
    
    # Add header bar
    header_shape = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(1))
    header_shape.fill.solid()
    header_shape.fill.fore_color.rgb = DARK_BLUE
    header_shape.line.color.rgb = DARK_BLUE
    
    # Add title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(0.7))
    title_frame = title_box.text_frame
    p = title_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = WHITE
    
    # Add content
    content_box = slide.shapes.add_textbox(Inches(0.7), Inches(1.5), Inches(8.6), Inches(5.5))
    text_frame = content_box.text_frame
    text_frame.word_wrap = True
    
    for i, item in enumerate(content_items):
        if i == 0:
            p = text_frame.paragraphs[0]
        else:
            p = text_frame.add_paragraph()
        
        p.text = item
        p.font.size = Pt(20)
        p.font.color.rgb = DARK_GRAY
        p.space_before = Pt(6)
        p.space_after = Pt(6)
        p.level = 0
        
        # Add bullet point
        if item.startswith("•"):
            p.text = item[1:].strip()
    
    return slide

# Slide 1: Welcome/Title Slide
add_title_slide(prs, "Food Order Application", "Complete MERN Stack Solution")

# Slide 2: Index/Agenda
add_content_slide(prs, "Presentation Index", [
    "• Welcome & Introduction",
    "• Project Objectives",
    "• Problem Statement",
    "• Proposed Solution",
    "• Technical Stack",
    "• System Architecture",
    "• Key Features",
    "• Database Design",
    "• Deployment & Implementation",
    "• Future Enhancements & Conclusion"
])

# Slide 3: Objectives
add_content_slide(prs, "Project Objectives", [
    "• Build a complete, production-ready food ordering system",
    "• Provide seamless shopping experience for customers",
    "• Enable efficient food inventory management for employees",
    "• Implement role-based access control (customer, employee, admin)",
    "• Ensure secure user authentication using JWT tokens",
    "• Create responsive, user-friendly interfaces",
    "• Enable scalable backend infrastructure using MongoDB"
])

# Slide 4: Problem Statement
add_content_slide(prs, "Problem Statement", [
    "• Limited digital ordering solutions in food service industry",
    "• Manual inventory management is time-consuming and error-prone",
    "• No centralized platform for order tracking and management",
    "• Difficulty in managing multiple user roles (customers, staff, admin)",
    "• Lack of real-time order status updates",
    "• Need for mobile-responsive ordering interface",
    "• Scalability issues with monolithic systems"
])

# Slide 5: Proposed Solution
add_content_slide(prs, "Proposed Solution", [
    "• Full-stack MERN (MongoDB, Express, React, Node.js) application",
    "• Multi-user portal system with role-based dashboards",
    "• Real-time order management and tracking system",
    "• JWT-based authentication for secure user sessions",
    "• Responsive React frontend for all device types",
    "• RESTful API backend with comprehensive documentation",
    "• MongoDB database for flexible data structure and scalability"
])

# Slide 6: Technical Stack
add_content_slide(prs, "Technical Stack", [
    "Frontend:",
    "  └ React.js - UI framework with component-based architecture",
    "  └ React Router - Client-side navigation",
    "  └ Axios - HTTP client for API communication",
    "",
    "Backend:",
    "  └ Node.js + Express.js - Server runtime and framework",
    "  └ MongoDB + Mongoose - NoSQL database with ODM",
    "  └ JWT - Secure token-based authentication"
])

# Slide 7: System Architecture
add_content_slide(prs, "System Architecture", [
    "Three-Tier Architecture:",
    "",
    "Presentation Layer (Frontend):",
    "  └ React components, context API for state management",
    "",
    "Business Logic Layer (Backend):",
    "  └ Express routes, controllers, and middleware",
    "",
    "Data Layer (Database):",
    "  └ MongoDB collections: Users, Foods, Orders, Employees, Admins"
])

# Slide 8: Key Features
add_content_slide(prs, "Key Features", [
    "Customer Features:",
    "  └ Browse food items by category",
    "  └ Add/remove items from cart with quantity control",
    "  └ Place orders with real-time order tracking",
    "",
    "Employee Features:",
    "  └ Manage food inventory (add, edit, delete items)",
    "  └ Toggle food availability status",
    "",
    "Admin Features:",
    "  └ Manage employee accounts and permissions"
])

# Slide 9: Database Design
add_content_slide(prs, "Database Design & Models", [
    "Collections Implemented:",
    "",
    "• User Collection - Customer profiles with authentication",
    "• Food Collection - Menu items with pricing and categories",
    "• Order Collection - Order history and status tracking",
    "• Employee Collection - Employee accounts and permissions",
    "• Admin Collection - Admin accounts and system controls",
    "",
    "All secured with JWT-based authentication and role-based access control"
])

# Slide 10: Conclusion & Future
add_content_slide(prs, "Implementation & Future Enhancements", [
    "Current Implementation:",
    "  └ Complete backend API with all CRUD operations",
    "  └ Fully functional React frontend with responsive design",
    "",
    "Future Enhancements:",
    "  └ Payment gateway integration (Stripe, PayPal)",
    "  └ Real-time notifications & push alerts",
    "  └ Advanced analytics and reporting dashboards",
    "  └ Mobile app development (React Native)",
    "  └ AI-based food recommendations"
])

# Save presentation
prs.save(r'd:\program\node js\food_order\Food_Order_App_Presentation.pptx')
print("✅ PowerPoint presentation created successfully!")
print("📍 Location: d:\\program\\node js\\food_order\\Food_Order_App_Presentation.pptx")
